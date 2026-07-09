#!/usr/bin/env python
"""Render markdown report text to md/pptx/docx/pdf (graceful degradation)."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any, Literal

ReportFormat = Literal["md", "pptx", "docx", "pdf"]
MERMAID_BLOCK_RE = re.compile(r"```mermaid[ \t]*\r?\n(?P<body>.*?)(?:\r?\n```|$)", re.IGNORECASE | re.DOTALL)
MERMAID_ALLOWED_STARTERS = {"flowchart", "sequenceDiagram", "classDiagram", "stateDiagram-v2"}
MERMAID_RESERVED_IDS = {"end", "graph", "subgraph"}
MERMAID_SEQUENCE_RESERVED_IDS = {
    "actor",
    "activate",
    "alt",
    "and",
    "autonumber",
    "box",
    "break",
    "critical",
    "create",
    "deactivate",
    "destroy",
    "else",
    "end",
    "loop",
    "note",
    "opt",
    "par",
    "participant",
    "rect",
}
MERMAID_FORBIDDEN_DIRECTIVES = ("style ", "classDef ", "click ", "class ")
MERMAID_SEQUENCE_DECL_RE = re.compile(
    r"^\s*(?:participant|actor)\s+(?P<actor>[A-Za-z_][\w.-]*)(?:\s+as\s+(?P<alias>.+?))?\s*$",
    re.IGNORECASE,
)
MERMAID_SEQUENCE_MESSAGE_RE = re.compile(
    r"^\s*(?P<from>[A-Za-z_][\w.-]*)\s*[-=]+[)>xX-]*\+?\s*(?P<to>[A-Za-z_][\w.-]*)\s*:",
)


def _line_number(text: str, index: int) -> int:
    return text.count("\n", 0, index) + 1


def _mermaid_issue(
    *,
    block: int,
    line: int,
    severity: str,
    code: str,
    message: str,
    snippet: str = "",
) -> dict[str, Any]:
    issue: dict[str, Any] = {
        "block": block,
        "line": line,
        "severity": severity,
        "code": code,
        "message": message,
    }
    if snippet:
        issue["snippet"] = snippet[:160]
    return issue


def _first_mermaid_statement(lines: list[str]) -> tuple[int, str] | None:
    for idx, line in enumerate(lines, start=1):
        stripped = line.strip()
        if not stripped or stripped.startswith("%%"):
            continue
        return idx, stripped
    return None


def _looks_like_raw_path_id(line: str) -> bool:
    stripped = line.strip()
    if not stripped or stripped.startswith("%%"):
        return False
    first = re.split(r"\s+", stripped, maxsplit=1)[0]
    if first.startswith(("http://", "https://")):
        return False
    return bool(re.search(r"(?:^|[A-Za-z]:)[\w.-]+[\\/][\w./\\-]+$", first))


def _is_quoted_mermaid_label(value: str) -> bool:
    stripped = value.strip()
    return len(stripped) >= 2 and stripped.startswith('"') and stripped.endswith('"')


def _sequence_alias_needs_quotes(alias: str) -> bool:
    return bool(re.search(r"[()/,:]", alias))


def _add_sequence_actor_issue(
    errors: list[dict[str, Any]],
    *,
    block_index: int,
    line: int,
    actor_id: str,
    snippet: str,
) -> None:
    if actor_id.lower() not in MERMAID_SEQUENCE_RESERVED_IDS:
        return
    errors.append(
        _mermaid_issue(
            block=block_index,
            line=line,
            severity="error",
            code="reserved_sequence_actor_id",
            message="Do not use Mermaid sequence keywords as participant or actor IDs; use short IDs like CinePart or TargetActor.",
            snippet=snippet,
        )
    )


def _validate_mermaid_block(body: str, *, block_index: int, start_line: int) -> dict[str, Any]:
    lines = body.splitlines()
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    first_statement = _first_mermaid_statement(lines)
    diagram_type = ""

    if first_statement is None:
        errors.append(
            _mermaid_issue(
                block=block_index,
                line=start_line,
                severity="error",
                code="empty_mermaid_block",
                message="Mermaid block is empty.",
            )
        )
    else:
        first_idx, first_line = first_statement
        diagram_type = first_line.split()[0]
        if diagram_type not in MERMAID_ALLOWED_STARTERS:
            errors.append(
                _mermaid_issue(
                    block=block_index,
                    line=start_line + first_idx - 1,
                    severity="error",
                    code="unsupported_diagram_type",
                    message=(
                        "Use one of: flowchart, sequenceDiagram, classDiagram, "
                        "or stateDiagram-v2."
                    ),
                    snippet=first_line,
                )
            )

    for idx, line in enumerate(lines, start=1):
        stripped = line.strip()
        if not stripped or stripped.startswith("%%"):
            continue
        absolute_line = start_line + idx - 1
        lowered = stripped.lower()
        if any(lowered.startswith(directive.lower()) for directive in MERMAID_FORBIDDEN_DIRECTIVES):
            errors.append(
                _mermaid_issue(
                    block=block_index,
                    line=absolute_line,
                    severity="error",
                    code="forbidden_directive",
                    message="Avoid style/class/click directives; renderer theme and security settings should control presentation.",
                    snippet=stripped,
                )
            )
        if re.search(r"\[[^\]\"']*[():,][^\]]*\]", stripped):
            errors.append(
                _mermaid_issue(
                    block=block_index,
                    line=absolute_line,
                    severity="error",
                    code="unquoted_special_label",
                    message="Wrap labels containing parentheses, brackets, commas, or colons in double quotes.",
                    snippet=stripped,
                )
            )
        reserved_match = re.match(r"^(end|graph|subgraph)(?:\b|\[|\(|\{|--|---|==|-->)", stripped, re.IGNORECASE)
        if reserved_match and reserved_match.group(1).lower() in MERMAID_RESERVED_IDS:
            errors.append(
                _mermaid_issue(
                    block=block_index,
                    line=absolute_line,
                    severity="error",
                    code="reserved_node_id",
                    message="Do not use Mermaid reserved keywords as node IDs.",
                    snippet=stripped,
                )
            )
        if _looks_like_raw_path_id(stripped):
            warnings.append(
                _mermaid_issue(
                    block=block_index,
                    line=absolute_line,
                    severity="warning",
                    code="raw_path_node_id",
                    message="Do not use raw file paths as node IDs; put paths in quoted labels or nearby text.",
                    snippet=stripped,
                )
            )
        if diagram_type == "sequenceDiagram":
            decl_match = MERMAID_SEQUENCE_DECL_RE.match(stripped)
            if decl_match:
                actor_id = decl_match.group("actor")
                alias = decl_match.group("alias") or ""
                _add_sequence_actor_issue(
                    errors,
                    block_index=block_index,
                    line=absolute_line,
                    actor_id=actor_id,
                    snippet=stripped,
                )
                if alias and _sequence_alias_needs_quotes(alias) and not _is_quoted_mermaid_label(alias):
                    errors.append(
                        _mermaid_issue(
                            block=block_index,
                            line=absolute_line,
                            severity="error",
                            code="unquoted_sequence_alias",
                            message="Quote sequence participant aliases that contain parentheses, slashes, commas, or colons.",
                            snippet=stripped,
                        )
                    )
            message_match = MERMAID_SEQUENCE_MESSAGE_RE.match(stripped)
            if message_match:
                for actor_id in (message_match.group("from"), message_match.group("to")):
                    _add_sequence_actor_issue(
                        errors,
                        block_index=block_index,
                        line=absolute_line,
                        actor_id=actor_id,
                        snippet=stripped,
                    )

    return {
        "index": block_index,
        "startLine": start_line,
        "diagramType": diagram_type,
        "ok": not errors,
        "errorCount": len(errors),
        "warningCount": len(warnings),
        "errors": errors,
        "warnings": warnings,
    }


def validate_mermaid_diagrams(text: str) -> dict[str, Any]:
    """Extract Mermaid fences and flag CI-safe render risks without rendering."""
    blocks: list[dict[str, Any]] = []
    for index, match in enumerate(MERMAID_BLOCK_RE.finditer(text or ""), start=1):
        body = match.group("body")
        start_line = _line_number(text, match.start("body"))
        blocks.append(_validate_mermaid_block(body, block_index=index, start_line=start_line))
    error_count = sum(block["errorCount"] for block in blocks)
    warning_count = sum(block["warningCount"] for block in blocks)
    return {
        "ok": error_count == 0,
        "blockCount": len(blocks),
        "errorCount": error_count,
        "warningCount": warning_count,
        "blocks": blocks,
    }


def render_report(
    text: str,
    *,
    format: ReportFormat = "md",
    output_path: str | Path | None = None,
) -> dict[str, Any]:
    """Render report text. Markdown always succeeds; other formats degrade if deps missing."""
    fmt = str(format or "md").strip().lower()
    if fmt not in {"md", "pptx", "docx", "pdf"}:
        fmt = "md"

    result: dict[str, Any] = {
        "ok": True,
        "format": fmt,
        "requestedFormat": format,
        "degraded": False,
        "outputPath": "",
        "notes": [],
    }
    mermaid_validation = validate_mermaid_diagrams(text)
    if mermaid_validation["blockCount"] > 0:
        result["mermaidValidation"] = mermaid_validation

    if output_path is None:
        suffix = {"md": ".md", "pptx": ".pptx", "docx": ".docx", "pdf": ".pdf"}[fmt]
        output_path = Path("report") / f"report{suffix}"
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    if fmt == "md":
        out.write_text(text, encoding="utf-8")
        result["outputPath"] = str(out.resolve())
        return result

    try:
        if fmt == "docx":
            _render_docx(text, out)
        elif fmt == "pptx":
            _render_pptx(text, out)
        elif fmt == "pdf":
            _render_pdf(text, out)
        result["outputPath"] = str(out.resolve())
        return result
    except ImportError as exc:
        md_fallback = out.with_suffix(".md")
        md_fallback.write_text(text, encoding="utf-8")
        result["degraded"] = True
        result["format"] = "md"
        result["outputPath"] = str(md_fallback.resolve())
        result["notes"].append(f"Missing dependency for {fmt}: {exc}. Wrote UTF-8 markdown instead.")
        return result
    except Exception as exc:
        md_fallback = out.with_suffix(".md")
        md_fallback.write_text(text, encoding="utf-8")
        result["degraded"] = True
        result["format"] = "md"
        result["outputPath"] = str(md_fallback.resolve())
        result["notes"].append(f"Render failed for {fmt}: {exc}. Wrote UTF-8 markdown instead.")
        return result


def _split_sections(text: str) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    current_title = "Report"
    current_lines: list[str] = []
    for line in text.splitlines():
        if line.startswith("#"):
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
            current_title = line.lstrip("#").strip() or "Section"
            current_lines = []
        else:
            current_lines.append(line)
    if current_lines or not sections:
        sections.append((current_title, "\n".join(current_lines).strip()))
    return sections


def _render_docx(text: str, out: Path) -> None:
    from docx import Document  # type: ignore[import-untyped]

    doc = Document()
    for title, body in _split_sections(text):
        doc.add_heading(title, level=1)
        for paragraph in body.split("\n\n"):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
    doc.save(out)


def _render_pptx(text: str, out: Path) -> None:
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    prs = Presentation()
    for title, body in _split_sections(text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for index, paragraph in enumerate(body.split("\n")):
            if not paragraph.strip():
                continue
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.text = paragraph.strip()
            p.font.size = Pt(16)
    prs.save(out)


def _render_pdf(text: str, out: Path) -> None:
    from reportlab.lib.pagesizes import letter  # type: ignore[import-untyped]
    from reportlab.pdfgen import canvas  # type: ignore[import-untyped]

    c = canvas.Canvas(str(out), pagesize=letter)
    width, height = letter
    y = height - 72
    for line in text.splitlines():
        if y < 72:
            c.showPage()
            y = height - 72
        c.drawString(72, y, line[:120])
        y -= 14
    c.save()


def main() -> int:
    parser = argparse.ArgumentParser(description="Render a markdown report to md/pptx/docx/pdf.")
    parser.add_argument("--input", default="-")
    parser.add_argument("--format", default="md", choices=["md", "pptx", "docx", "pdf"])
    parser.add_argument("--output", default="")
    args = parser.parse_args()
    if args.input == "-":
        import sys

        text = sys.stdin.read()
    else:
        text = Path(args.input).read_text(encoding="utf-8", errors="replace")
    payload = render_report(text, format=args.format, output_path=args.output or None)
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0 if payload.get("ok") else 1


if __name__ == "__main__":
    raise SystemExit(main())
